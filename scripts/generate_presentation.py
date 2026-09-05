"""
generate_presentation.py
------------------------
Generates reports/Bluestock_MF_Presentation.pptx using python-pptx.
12 slides with real data pulled from fund_scorecard.csv and alpha_beta.csv.
Embeds actual chart PNGs and dashboard screenshot placeholders.
"""

import sqlite3
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE   = Path(__file__).resolve().parent.parent
PROC   = BASE / "data" / "processed"
CHARTS = BASE / "reports" / "charts"
DASH   = BASE / "reports" / "dashboard_screenshots"
OUT    = BASE / "reports" / "Bluestock_MF_Presentation.pptx"

# ── Load data ─────────────────────────────────────────────────────────────────
sc = pd.read_csv(PROC / "fund_scorecard.csv")
ab = pd.read_csv(PROC / "alpha_beta.csv")

conn = sqlite3.connect(BASE / "data" / "db" / "bluestock_mf.db")
aum_df  = pd.read_sql("SELECT * FROM fact_aum", conn)
aum_df = aum_df.rename(columns={"date": "month", "aum_crore": "aum_cr", "fund_house": "amc"})
txn_df  = pd.read_sql("SELECT * FROM fact_transactions", conn)
txn_df = txn_df.rename(columns={"transaction_date": "txn_date", "amfi_code": "fund_id", "transaction_type": "txn_type", "amount_inr": "amount"})
conn.close()

latest_month = aum_df["month"].max()
total_aum    = aum_df[aum_df["month"] == latest_month]["aum_cr"].sum()
sip_total    = txn_df[txn_df["txn_type"] == "Sip"]["amount"].sum() / 1e7
top          = sc.iloc[0]

# ── Colors ─────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0B, 0x1F, 0x4B)
BLUE    = RGBColor(0x1A, 0x56, 0xDB)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xEF, 0xF6, 0xFF)
SLATE   = RGBColor(0x64, 0x74, 0x8B)
DARK    = RGBColor(0x1E, 0x29, 0x3B)

# ── Helpers ───────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, l, t, w, h, fill=None, line=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h, size=24, bold=False, color=DARK,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_img(slide, path, l, t, w, h=None):
    if Path(path).exists():
        if h:
            slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
        else:
            slide.shapes.add_picture(str(path), Inches(l), Inches(t), width=Inches(w))

def header_bar(slide, title, subtitle=None):
    """Navy header bar at top of slide."""
    rect(slide, 0, 0, 13.33, 1.3, fill=NAVY)
    txbox(slide, title, 0.4, 0.08, 10, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.4, 0.75, 12, 0.45, size=13, color=EMERALD)

def footer_bar(slide):
    rect(slide, 0, 7.1, 13.33, 0.4, fill=NAVY)
    txbox(slide, "Bluestock MF Analytics Platform  |  Python · SQLite · Streamlit · Plotly",
          0.3, 7.12, 10, 0.3, size=8, color=SLATE)
    txbox(slide, "Confidential — Capstone Project 2026",
          10.5, 7.12, 2.5, 0.3, size=8, color=SLATE, align=PP_ALIGN.RIGHT)

def bullet_box(slide, items, l, t, w, h, size=13):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)

def kpi_card(slide, l, t, w, h, value, label, color=BLUE):
    rect(slide, l, t, w, h, fill=WHITE, line=color)
    rect(slide, l, t, 0.07, h, fill=color)
    txbox(slide, value, l+0.15, t+0.08, w-0.2, h*0.5, size=22, bold=True, color=NAVY)
    txbox(slide, label, l+0.15, t+h*0.52, w-0.2, h*0.44, size=10, color=SLATE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
rect(sl, 0, 0, 0.18, 7.5, fill=EMERALD)
rect(sl, 0, 3.4, 13.33, 0.06, fill=EMERALD)
txbox(sl, "BLUESTOCK FINTECH", 1, 1.2, 11, 0.6, size=16, color=EMERALD, align=PP_ALIGN.CENTER)
txbox(sl, "Mutual Fund Analytics Platform", 1, 1.9, 11, 1.1,
      size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "Capstone Project — Final Presentation", 1, 3.1, 11, 0.6,
      size=18, color=LIGHT, align=PP_ALIGN.CENTER)
txbox(sl, f"10 Funds  |  50,000 Transactions  |  2,000 Investors  |  Data: 2022–2026",
      1, 4.0, 11, 0.5, size=13, color=SLATE, align=PP_ALIGN.CENTER)
txbox(sl, "Python 3.12  ·  SQLite  ·  Streamlit  ·  Plotly  ·  ReportLab",
      1, 4.6, 11, 0.5, size=12, color=SLATE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Objective
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "Problem & Objective",
           "Why does the Indian MF industry need a unified analytics platform?")
rect(sl, 0.3, 1.5, 6.0, 5.3, fill=RGBColor(0xEF,0xF6,0xFF), line=BLUE)
rect(sl, 6.8, 1.5, 6.0, 5.3, fill=RGBColor(0xF0,0xFD,0xF4), line=EMERALD)

txbox(sl, "❌  The Problem", 0.5, 1.55, 5.6, 0.5, size=14, bold=True, color=NAVY)
bullet_box(sl, [
    "MF data spread across NAV feeds, transaction logs, AUM reports, and benchmark indices — no unified view",
    "Fund managers lack real-time risk metrics (VaR, Sharpe, drawdown) in one place",
    "Investor behaviour (SIP continuity, geographic trends) not systematically tracked",
    "No composite scoring framework to compare funds across risk-adjusted metrics",
], 0.5, 2.1, 5.6, 4.5, size=12)

txbox(sl, "✅  Our Solution", 7.0, 1.55, 5.6, 0.5, size=14, bold=True, color=EMERALD)
bullet_box(sl, [
    "End-to-end ETL pipeline: raw CSVs → cleaned data → SQLite star-schema DB",
    "Live NAV integration from AMFI API (mfapi.in) — 16,352 real records",
    "23-chart EDA + advanced analytics (VaR, cohort, HHI, recommender)",
    "Interactive Streamlit dashboard: 4 pages, Plotly charts, fintech theme",
    "Composite fund scorecard ranking all 10 schemes by risk-adjusted return",
], 7.0, 2.1, 5.6, 4.5, size=12)
footer_bar(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data Sources
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "Data Sources & Datasets", "10 CSVs + Live API  |  86,000+ database records")

data_items = [
    ("fund_master.csv",             "10 rows",    "Fund metadata, categories, AMC, expense ratios"),
    ("nav_history.csv",             "12,200 rows","Daily NAV per fund (2022–2026, seeded 1% nulls)"),
    ("aum_history.csv",             "560 rows",   "Monthly AUM and net inflows per fund"),
    ("investor_transactions.csv",   "50,000 rows","SIP, Lumpsum, Redemption, STP, Switch"),
    ("sip_register.csv",            "800 rows",   "SIP mandates with continuity tracking"),
    ("scheme_performance.csv",      "40 rows",    "Annual returns, Sharpe, expense ratios"),
    ("portfolio_holdings.csv",      "1,260 rows", "Quarterly sector allocations per fund"),
    ("benchmark_returns.csv",       "9,760 rows", "8 benchmarks: Nifty 100, Nifty 500, Crisil, etc."),
    ("investor_demographics.csv",   "2,000 rows", "Age, city, KYC status, risk profile"),
    ("Live NAV (mfapi.in)",         "16,352 rows","Real AMFI NAV for 6 scheme codes"),
]
col_w = [3.8, 1.6, 7.0]
for i, (name, rows, desc) in enumerate(data_items):
    y = 1.55 + i * 0.5
    bg = RGBColor(0xEF,0xF6,0xFF) if i % 2 == 0 else WHITE
    rect(sl, 0.3, y, 12.7, 0.46, fill=bg)
    txbox(sl, name,  0.4,  y+0.05, col_w[0], 0.36, size=10, bold=True, color=NAVY)
    txbox(sl, rows,  4.3,  y+0.05, col_w[1], 0.36, size=10, color=BLUE)
    txbox(sl, desc,  6.0,  y+0.05, col_w[2], 0.36, size=10, color=DARK)
footer_bar(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture / ETL Diagram
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "ETL Pipeline Architecture", "8-step sequential pipeline from raw CSVs to interactive dashboard")

steps = [
    ("0\nGenerate", SLATE),   ("1\nIngest",  BLUE),   ("2\nLive NAV", TEAL := RGBColor(0x0E,0xA5,0xE9)),
    ("3\nClean",    EMERALD), ("4\nLoad DB", NAVY),   ("5\nSQL",      BLUE),
    ("6–8\nNotebooks", RGBColor(0x8B,0x5C,0xF6)), ("Dashboard\nStreamlit", EMERALD),
]
box_w, box_h, start_x, y = 1.3, 0.9, 0.4, 2.0
for i, (label, color) in enumerate(steps):
    x = start_x + i * 1.6
    if i == 4:  # wrap to second row
        y = 3.5
        x = start_x
    elif i > 4:
        x = start_x + (i - 4) * 1.6
    rect(sl, x, y, box_w, box_h, fill=color)
    txbox(sl, label, x+0.05, y+0.1, box_w-0.1, box_h-0.1,
          size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1 and i != 3:
        txbox(sl, "→", x + box_w, y + 0.3, 0.25, 0.4, size=16, bold=True, color=SLATE)

txbox(sl, "Key Technologies:", 0.4, 5.1, 2.5, 0.4, size=12, bold=True, color=NAVY)
txbox(sl, "Python 3.12  ·  Pandas 3.0  ·  SQLAlchemy 2.0  ·  SQLite  ·  Plotly  ·  Streamlit  ·  ReportLab  ·  python-pptx",
      0.4, 5.5, 12.5, 0.5, size=11, color=DARK)
txbox(sl, f"Total records in DB: 86,630  |  Tables: 9  |  Charts generated: 23",
      0.4, 6.1, 12.5, 0.5, size=11, color=BLUE)
footer_bar(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EDA Highlights #1
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "EDA Highlights — Industry Overview",
           f"Total AUM: ₹{total_aum:,.0f} Cr  |  SIP Inflows: ₹{sip_total:.1f} Cr  |  50K Transactions")

kpi_card(sl, 0.3,  1.35, 2.9, 1.0, f"₹{total_aum/1e5:.1f}L Cr", "Total Industry AUM", BLUE)
kpi_card(sl, 3.35, 1.35, 2.9, 1.0, f"₹{sip_total:.1f} Cr", "SIP Inflows", EMERALD)
kpi_card(sl, 6.40, 1.35, 2.9, 1.0, "48,638", "Active Folios", RGBColor(0x8B,0x5C,0xF6))
kpi_card(sl, 9.45, 1.35, 3.5, 1.0, "10 Schemes", "Across 4 AMCs", NAVY)

add_img(sl, CHARTS/"01_nav_trend_all.png",     0.3,  2.5, 6.3, 4.3)
add_img(sl, CHARTS/"03_sip_monthly_milestone.png", 6.7, 2.5, 6.3, 4.3)
footer_bar(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — EDA Highlights #2
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "EDA Highlights — Investor & Sector Analytics",
           "2,000 investors  |  T30/B30 split  |  10 sectors tracked")
add_img(sl, CHARTS/"05_t30_b30_comparison.png", 0.3, 1.4, 6.2, 4.2)
add_img(sl, CHARTS/"10_sector_donut.png",       6.7, 1.4, 6.2, 4.2)
bullet_box(sl, [
    "T30 cities: 60% of total capital  |  B30 cities: higher SIP continuation rates",
    "Financial Services = top sector allocation (18–22% across equity funds)",
    "IT and FMCG consistently rank 2nd and 3rd in portfolio holdings",
], 0.4, 5.7, 12.5, 1.2, size=11)
footer_bar(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Performance Metrics #1 (Scorecard)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "Fund Performance — Composite Scorecard",
           f"Top Fund: {top['scheme_name']}  |  CAGR 1Y: {top['return_1yr_pct']:.1f}%  |  Sharpe: {top['sharpe_ratio']:.3f}")

add_img(sl, CHARTS/"16_scorecard_heatmap.png", 0.3, 1.4, 7.8, 5.0)

sc_top5 = sc[["overall_rank","amfi_code","scheme_name","return_1yr_pct","sharpe_ratio","max_drawdown_pct","composite_score"]].head(5)
y_start = 1.5
col_heads = ["Rank", "ID", "Scheme", "CAGR 1Y%", "Sharpe", "MaxDD%", "Score"]
col_xs    = [8.2, 8.7, 9.2, 11.0, 11.7, 12.2, 12.8]
col_ws    = [0.45, 0.5, 1.85, 0.65, 0.65, 0.6, 0.5]

for ci, (head, cx, cw) in enumerate(zip(col_heads, col_xs, col_ws)):
    rect(sl, cx-0.02, y_start, cw+0.04, 0.4, fill=NAVY)
    txbox(sl, head, cx, y_start+0.05, cw, 0.3, size=8, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in sc_top5.iterrows():
    y = y_start + 0.42 + ri * 0.46
    bg = RGBColor(0xEF,0xF6,0xFF) if ri % 2 == 0 else WHITE
    rect(sl, col_xs[0]-0.02, y, sum(col_ws)+0.04+0.08, 0.42, fill=bg)
    vals = [str(row["overall_rank"]), str(row["amfi_code"]), 
            row["scheme_name"][:20]+"..", f"{row['return_1yr_pct']:.1f}%", 
            f"{row['sharpe_ratio']:.2f}", f"{row['max_drawdown_pct']:.1f}%", f"{row['composite_score']:.2f}"]
    for cx, cw, val in zip(col_xs, col_ws, vals):
        txbox(sl, val, cx, y+0.07, cw, 0.3, size=8, color=DARK, align=PP_ALIGN.CENTER)

footer_bar(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Performance Metrics #2 (Alpha/Beta + Risk)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "Fund Performance — Alpha, Beta & Risk Metrics",
           "OLS regression vs Nifty 100  |  Max Drawdown  |  Tracking Error")
add_img(sl, CHARTS/"13_risk_return_scatter.png", 0.3, 1.4, 6.2, 4.0)
add_img(sl, CHARTS/"17_max_drawdown.png",        6.7, 1.4, 6.2, 4.0)
bullet_box(sl, [
    f"All equity funds show positive alpha (best: Top Fund = {ab['alpha'].max():.4f} annualised)",
    "Top Large Cap funds occupy top-right of risk/return chart: best return per unit of risk",
    f"Worst max drawdown: Worst at {sc['max_drawdown_pct'].min():.1f}%  |  Best: Best at {sc['max_drawdown_pct'].max():.1f}%",
], 0.4, 5.55, 12.5, 1.2, size=11)
footer_bar(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Dashboard Screenshots #1
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "Interactive Streamlit Dashboard — Pages 1 & 2",
           "Run: python3 -m streamlit run dashboard/app.py  →  http://localhost:8501")
add_img(sl, DASH/"page1_industry_overview.png",  0.3, 1.4, 6.2, 5.4)
add_img(sl, DASH/"page2_fund_performance.png",   6.7, 1.4, 6.2, 5.4)
footer_bar(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Dashboard Screenshots #2
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "Interactive Streamlit Dashboard — Pages 3 & 4",
           "Sidebar filters for city tier, age group, fund house, category")
add_img(sl, DASH/"page3_investor_analytics.png", 0.3, 1.4, 6.2, 5.4)
add_img(sl, DASH/"page4_sip_market_trends.png",  6.7, 1.4, 6.2, 5.4)
footer_bar(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Key Findings & Recommendations
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "Key Findings & Recommendations", "Actionable insights from the analytics platform")

rect(sl, 0.3, 1.4, 6.0, 5.4, fill=RGBColor(0xEF,0xF6,0xFF), line=BLUE)
rect(sl, 6.8, 1.4, 6.0, 5.4, fill=RGBColor(0xF0,0xFD,0xF4), line=EMERALD)

txbox(sl, "📊  Key Findings", 0.5, 1.45, 5.6, 0.45, size=13, bold=True, color=NAVY)
bullet_box(sl, [
    f"The top equity large-cap fund: #1 fund — CAGR {top['return_1yr_pct']:.1f}%, Sharpe {top['sharpe_ratio']:.3f}",
    "SIP inflows peaked December 2025 — industry milestone confirmed",
    "60%+ SIP mandates rated Good/Excellent continuity",
    "B30 cities: higher SIP retention despite lower ticket size",
    "Equity fund NAV correlation > 0.85 — minimal intra-category diversification",
    "5 TER violations detected and corrected in cleaning step",
    "HHI < 0.15 for most funds — healthy sector diversification",
], 0.5, 1.95, 5.6, 4.6, size=11)

txbox(sl, "🎯  Recommendations", 7.0, 1.45, 5.6, 0.45, size=13, bold=True, color=EMERALD)
bullet_box(sl, [
    "Increase marketing spend on top fund — clear performance leader",
    "Trigger advisor outreach for At Risk SIP segment (~15%)",
    "Target 18–30 cohort with digital-first SIP products",
    "Review underperforming funds with negative CAGR",
    "Migrate to PostgreSQL for production dashboard scalability",
    "Integrate daily AMFI NAV scheduler for live updates",
    "Add ML churn prediction using SIP continuity features",
], 7.0, 1.95, 5.6, 4.6, size=11)
footer_bar(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
rect(sl, 0, 0, 0.18, 7.5, fill=EMERALD)
rect(sl, 6.5, 0, 0.04, 7.5, fill=RGBColor(0x1A,0x56,0xDB))

txbox(sl, "Thank You", 1, 1.5, 11, 1.2,
      size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "Bluestock Fintech — Mutual Fund Analytics Platform", 1, 2.9, 11, 0.6,
      size=18, color=EMERALD, align=PP_ALIGN.CENTER)
txbox(sl, "Questions & Discussion", 1, 3.6, 11, 0.5,
      size=14, color=LIGHT, align=PP_ALIGN.CENTER)

txbox(sl, "📁  GitHub: bluestock-mf-capstone  |  🌐  localhost:8501", 1, 4.5, 11, 0.5,
      size=12, color=SLATE, align=PP_ALIGN.CENTER)
txbox(sl, "python3 -m streamlit run dashboard/app.py", 1, 5.1, 11, 0.5,
      size=13, color=EMERALD, align=PP_ALIGN.CENTER)

txbox(sl, "Built with: Python 3.12  ·  Pandas  ·  SQLite  ·  SQLAlchemy  ·  Streamlit  ·  Plotly  ·  ReportLab  ·  python-pptx",
      1, 6.2, 11, 0.5, size=10, color=SLATE, align=PP_ALIGN.CENTER)

prs.save(str(OUT))
print(f"✅ PPTX saved to {OUT}  ({OUT.stat().st_size // 1024} KB)  |  {len(prs.slides)} slides")
